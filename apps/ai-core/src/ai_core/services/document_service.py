"""
Document processing: chunking and embedding using OpenAI.
"""
from typing import List, Tuple, Dict, Any
from sqlalchemy.orm import Session
from shared.database.models import Document, KnowledgeChunk, KnowledgeBase, Tenant
from shared.utils.storage import write_metadata
from openai import OpenAI
import os
import hashlib
import random
import io
import csv
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from shared.vector.qdrant import qdrant_service
from shared.config.tuning import chunking
from shared.crypto.crypto_service import crypto_service
from shared.plans.registry import resolve_plan_label, get_plan

try:
    import tiktoken  # type: ignore
except Exception:
    tiktoken = None  # fallback
import logging
import re
import pandas as pd
import numpy as np


def chunk_text(text: str, chunk_size: int = 700, overlap: int = 100) -> List[str]:
    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = max(0, end - overlap)
    return chunks


class DocumentService:
    def __init__(self, db: Session):
        self.db = db
        # Initialize OpenAI client lazily but ensure attributes exist
        try:
            from shared.security.secret_manager import secret_manager

            api_key = secret_manager.get("OPENAI_API_KEY")
        except Exception:
            api_key = os.getenv("OPENAI_API_KEY")
        try:
            from openai import OpenAI as _OpenAI

            self.client = _OpenAI(api_key=api_key) if api_key else None
            self.openai_client = self.client if api_key else None
        except Exception:
            self.client = None
            self.openai_client = None

    # ------------------------------
    # Normalized connector ingestion
    # ------------------------------
    def process_normalized_records(self, records: list) -> int:
        """Ingest normalized connector records with dedup detection.

        Records are dict-like objects with fields per NormalizedRecord.
        """
        from shared.database.models import (
            KnowledgeBase,
            Document,
            KnowledgeChunk,
        )
        import uuid as _uuid

        count = 0
        for rec in records:
            trans = self.db.begin()
            try:
                tenant_id = (
                    rec.tenant_id if hasattr(rec, "tenant_id") else rec.get("tenant_id")
                )
                source_system = (
                    rec.source_system
                    if hasattr(rec, "source_system")
                    else rec.get("source_system")
                )
                external_id = (
                    rec.external_id
                    if hasattr(rec, "external_id")
                    else rec.get("external_id")
                )
                title = rec.title if hasattr(rec, "title") else rec.get("title")
                content = rec.content if hasattr(rec, "content") else rec.get("content")
                meta = {
                    "source_system": source_system,
                    "external_id": external_id,
                    "owner": getattr(rec, "owner", None)
                    if hasattr(rec, "owner")
                    else rec.get("owner"),
                    "classification": getattr(rec, "classification", None)
                    if hasattr(rec, "classification")
                    else rec.get("classification"),
                }
                # Find or create a default KB for connector ingestion
                kb = (
                    self.db.query(KnowledgeBase)
                    .filter(
                        KnowledgeBase.tenant_id == tenant_id,
                        KnowledgeBase.name == "connector",
                    )
                    .first()
                )
                if not kb:
                    kb = KnowledgeBase(tenant_id=tenant_id, name="connector")
                    self.db.add(kb)
                    self.db.commit()
                    self.db.refresh(kb)
                # Dedup: existing document with same external id and source
                existing = (
                    self.db.query(Document)
                    .filter(
                        Document.knowledge_base_id == kb.id,
                        Document.meta["external_id"].astext == external_id,
                        Document.meta["source_system"].astext == source_system,
                    )
                    .first()
                )
                if existing:
                    doc = existing
                else:
                    doc = Document(
                        knowledge_base_id=kb.id,
                        title=title,
                        content=content[:160],
                        meta=meta,
                        status="PROCESSING",
                    )
                    self.db.add(doc)
                    self.db.commit()
                    self.db.refresh(doc)
                # Create one chunk per record (can be extended to chunking later)
                kc = KnowledgeChunk(
                    id=_uuid.uuid4(),
                    document_id=doc.id,
                    content=content[:1600],
                    chunk_index=0,
                    embedding=None,
                    meta={"source": source_system},
                )
                # Embed connector chunk
                try:
                    vecs = self.embed([kc.content])
                    if vecs:
                        kc.embedding = vecs[0]
                except Exception:
                    pass
                self.db.add(kc)
                count += 1
                # Best-effort upsert to vector store
                try:
                    from shared.vector.qdrant import qdrant_service

                    qdrant_service.upsert_knowledge_chunks(
                        tenant_id,
                        [
                            {
                                "id": str(kc.id),
                                "embedding": kc.embedding,
                                "document_id": str(doc.id),
                                "document_title": doc.title,
                                "content": kc.content,
                                "chunk_index": kc.chunk_index,
                                "metadata": kc.meta or {},
                            }
                        ],
                    )
                except Exception:
                    pass
                trans.commit()
            except Exception:
                try:
                    trans.rollback()
                except Exception:
                    pass
                raise
        return count

    @staticmethod
    def _split_sentences(text: str) -> List[str]:
        # Lightweight sentence splitter using punctuation and line breaks
        # Avoid breaking on common abbreviations by a simple heuristic
        text = re.sub(r"\s+", " ", text)
        candidates = re.split(r"(?<=[\.!?])\s+(?=[A-Z(\[])", text)
        sentences: List[str] = []
        for s in candidates:
            s = s.strip()
            if not s:
                continue
            if len(s) < 3:
                continue
            sentences.append(s)
        return sentences

    def _build_chunks_with_metadata(
        self,
        text: str,
        target_chars: int = chunking.target_chars,
        overlap_sentences: int = chunking.sentence_overlap,
    ) -> List[Tuple[str, Dict[str, Any]]]:
        """Sentence-aware chunking with small overlap and chapter/page tagging.

        Recognizes page markers like [[PAGE:n]] if present.
        Detects chapter headings like "Chapter 3. ..." and tags subsequent chunks
        until the next heading.
        """
        # Detect simple page markers
        pages: List[Tuple[int, str]] = []
        page_matches = list(re.finditer(r"\[\[PAGE:(\d+)\]\]", text))
        if page_matches:
            last_idx = 0
            current_page = 1
            segments: List[Tuple[int, str]] = []
            for m in page_matches:
                seg = text[last_idx : m.start()]
                if seg.strip():
                    segments.append((current_page, seg))
                try:
                    current_page = int(m.group(1))
                except Exception:
                    pass
                last_idx = m.end()
            tail = text[last_idx:]
            if tail.strip():
                segments.append((current_page, tail))
            pages = segments
        else:
            pages = [(None, text)]  # type: ignore

        chunks: List[Tuple[str, Dict[str, Any]]] = []
        current_chapter_num: Any = None
        current_chapter_title: str = ""

        for page_num, page_text in pages:
            # Identify chapter heading at start of page or within first lines
            for line in page_text.splitlines()[:6]:
                m = re.match(
                    r"^\s*chapter\s+(\d+)\s*[\.:\-]?\s*(.*)$",
                    line.strip(),
                    flags=re.IGNORECASE,
                )
                if m:
                    try:
                        current_chapter_num = int(m.group(1))
                        current_chapter_title = (m.group(2) or "").strip()
                    except Exception:
                        pass
                    break

            sentences = self._split_sentences(page_text)
            if not sentences:
                continue

            if chunking.mode == "tokens" and tiktoken is not None:
                # Token-based windowing over sentences
                enc = tiktoken.get_encoding(
                    os.getenv("TIKTOKEN_ENCODING", "cl100k_base")
                )
                # Precompute tokenized sentences
                sent_tokens = [enc.encode(s) for s in sentences]
                buf_tokens: List[int] = []
                buf_sent_start = 0
                i = 0
                while i < len(sent_tokens):
                    stoks = sent_tokens[i]
                    if not buf_tokens:
                        buf_sent_start = i
                        buf_tokens = stoks[:]
                    else:
                        prospective_len = len(buf_tokens) + len(stoks) + 1
                        if prospective_len <= chunking.target_tokens:
                            buf_tokens.extend(stoks)
                        else:
                            # emit chunk
                            text_chunk = enc.decode(buf_tokens)
                            meta: Dict[str, Any] = {}
                            if page_num is not None:
                                meta["page"] = page_num
                            if current_chapter_num is not None:
                                meta["chapter_num"] = current_chapter_num
                            if current_chapter_title:
                                meta["chapter_title"] = current_chapter_title
                            chunks.append((text_chunk.strip(), meta))
                            # overlap in tokens by reusing last N tokens from previous buffer
                            overlap_tok = max(0, chunking.overlap_tokens)
                            if overlap_tok > 0:
                                buf_tokens = buf_tokens[-overlap_tok:] + stoks
                            else:
                                buf_tokens = stoks[:]
                            buf_sent_start = max(buf_sent_start, i - 1)
                    i += 1
                if buf_tokens:
                    text_chunk = enc.decode(buf_tokens)
                    meta: Dict[str, Any] = {}
                    if page_num is not None:
                        meta["page"] = page_num
                    if current_chapter_num is not None:
                        meta["chapter_num"] = current_chapter_num
                    if current_chapter_title:
                        meta["chapter_title"] = current_chapter_title
                    chunks.append((text_chunk.strip(), meta))
            else:
                # Character-size sentence windowing (fallback)
                buf: List[str] = []
                for i, s in enumerate(sentences):
                    if not buf:
                        buf.append(s)
                    else:
                        prospective = (" ".join(buf) + " " + s).strip()
                        if len(prospective) <= target_chars:
                            buf.append(s)
                        else:
                            # Emit chunk
                            text_chunk = " ".join(buf).strip()
                            meta: Dict[str, Any] = {}
                            if page_num is not None:
                                meta["page"] = page_num
                            if current_chapter_num is not None:
                                meta["chapter_num"] = current_chapter_num
                            if current_chapter_title:
                                meta["chapter_title"] = current_chapter_title
                            chunks.append((text_chunk, meta))
                            # Start new buffer with overlap
                            overlap = sentences[max(0, i - overlap_sentences) : i]
                            buf = overlap + [s]

                if buf:
                    text_chunk = " ".join(buf).strip()
                    meta: Dict[str, Any] = {}
                    if page_num is not None:
                        meta["page"] = page_num
                    if current_chapter_num is not None:
                        meta["chapter_num"] = current_chapter_num
                    if current_chapter_title:
                        meta["chapter_title"] = current_chapter_title
                    chunks.append((text_chunk, meta))

        # Merge too-short chunks with previous where possible
        merged: List[Tuple[str, Dict[str, Any]]] = []
        for t, m in chunks:
            if merged and len(t) < chunking.min_chars:
                prev_t, prev_m = merged[-1]
                merged[-1] = (prev_t + " " + t, prev_m)
            else:
                merged.append((t, m))
        return merged

    def plan_ingest(self, title: str, sample_text: str) -> Dict[str, Any]:
        """Ask the LLM to propose an ingestion plan.

        Returns JSON with keys: chunk_mode (tokens|chars), target_tokens, overlap_tokens,
        heading_regex (optional), use_local_embeddings (bool).
        """
        default = {
            "chunk_mode": chunking.mode,
            "target_tokens": chunking.target_tokens,
            "overlap_tokens": chunking.overlap_tokens,
            "heading_regex": r"^chapter\s+(\d+)\s*[\.:\-]?\s*(.*)$",
            "use_local_embeddings": False,
        }
        if not self.openai_client:
            return default
        try:
            prompt = (
                "You are an ingestion planner. Based on the TITLE and SAMPLE, output a strict JSON object with keys: "
                "chunk_mode (tokens|chars), target_tokens (int), overlap_tokens (int), heading_regex (string or null), "
                "use_local_embeddings (true|false). Focus on keeping paragraphs/list items intact."
            )
            msg = [
                {"role": "system", "content": prompt},
                {
                    "role": "user",
                    "content": f"TITLE: {title}\nSAMPLE:\n{sample_text[:4000]}",
                },
            ]
            resp = self.openai_client.chat.completions.create(
                model=os.getenv(
                    "RAG_PLANNER_MODEL", os.getenv("RAG_CHAT_MODEL", "gpt-4o-mini")
                ),
                temperature=0,
                messages=msg,
            )
            import json as _json

            raw = (resp.choices[0].message.content or "").strip()
            plan = _json.loads(raw)
            if not isinstance(plan, dict):
                return default
            out = default.copy()
            out.update({k: plan.get(k, out[k]) for k in out.keys()})
            return out
        except Exception:
            return default

    @staticmethod
    def _extract_chapter_info(text: str) -> Dict[str, Any]:
        """Extract simple chapter metadata from a text chunk.

        Looks for patterns like:
        - "Chapter 3. Program and Project Management Roles and Responsibilities"
        - "Chapter 4: Governance"
        """
        try:
            import re

            lines = [l.strip() for l in text.splitlines() if l.strip()]
            for line in lines[:5]:  # inspect only early lines of the chunk
                m = re.match(
                    r"^chapter\s+(\d+)\s*[\.:\-]?\s*(.*)$", line, flags=re.IGNORECASE
                )
                if m:
                    num = int(m.group(1))
                    title = (m.group(2) or "").strip()
                    return {"chapter_num": num, "chapter_title": title}
        except Exception:
            pass
        return {}

    def embed(self, inputs: List[str], force_local: bool = False) -> List[List[float]]:
        # Deterministic local embedding path (used for tabular/robust ingestion)
        if force_local:
            vectors: List[List[float]] = []
            dim = 256
            for x in inputs:
                try:
                    s = x if isinstance(x, str) else str(x)
                    s = s[:5000]
                    if not s.strip():
                        s = " "
                except Exception:
                    s = " "
                h = hashlib.sha256(s.encode("utf-8")).digest()
                rnd = random.Random(h)
                vectors.append([rnd.uniform(-1.0, 1.0) for _ in range(dim)])
            return vectors
        api_key = os.getenv("OPENAI_API_KEY")
        if api_key:
            # Sanitize inputs for OpenAI API: non-empty strings only
            safe_inputs: List[str] = []
            for x in inputs:
                try:
                    s = x if isinstance(x, str) else str(x)
                    s = s[:5000]
                    if s and s.strip():
                        safe_inputs.append(s)
                except Exception:
                    continue
            if not safe_inputs:
                return []
            try:
                # Batch requests to respect OpenAI per-request token limits
                def estimate_tokens(text: str) -> int:
                    # Rough heuristic: 4 chars per token
                    return max(1, len(text) // 4)

                MAX_TOKENS_PER_REQUEST = (
                    100_000  # conservative to avoid 400s on large batches
                )
                embeddings: List[List[float]] = []
                batch: List[str] = []
                tokens_in_batch = 0
                for t in safe_inputs:
                    t_tokens = estimate_tokens(t)
                    if batch and tokens_in_batch + t_tokens > MAX_TOKENS_PER_REQUEST:
                        resp = self.client.embeddings.create(
                            model=os.getenv(
                                "RAG_EMBED_MODEL", "text-embedding-3-large"
                            ),
                            input=batch,
                        )
                        embeddings.extend([d.embedding for d in resp.data])
                        batch = []
                        tokens_in_batch = 0
                    batch.append(t)
                    tokens_in_batch += t_tokens

                if batch:
                    resp = self.client.embeddings.create(
                        model=os.getenv("RAG_EMBED_MODEL", "text-embedding-3-large"),
                        input=batch,
                    )
                    embeddings.extend([d.embedding for d in resp.data])

                return embeddings
            except Exception:
                # Global fallback to deterministic local embeddings for all inputs
                vectors: List[List[float]] = []
                dim = 256
                for text in safe_inputs:
                    h = hashlib.sha256(text.encode("utf-8")).digest()
                    rnd = random.Random(h)
                    vectors.append([rnd.uniform(-1.0, 1.0) for _ in range(dim)])
                return vectors
        # Fallback deterministic embedding (no external dependency)
        vectors: List[List[float]] = []
        dim = 256
        for text in inputs:
            h = hashlib.sha256(text.encode("utf-8")).digest()
            # Expand hash deterministically
            rnd = random.Random(h)
            vec = [rnd.uniform(-1.0, 1.0) for _ in range(dim)]
            vectors.append(vec)
        return vectors

    def process_and_store(
        self,
        tenant_id: str,
        title: str,
        content: str,
        knowledge_base_id: str,
        progress_job_id: str | None = None,
        doc_meta: Dict[str, Any] | None = None,
    ) -> Tuple[str, int]:
        # Wrap entire ingest in a single transaction for atomicity
        trans = self.db.begin()
        try:
            # Validate tenant_id is a valid UUID
            import uuid

            try:
                uuid.UUID(tenant_id)
            except ValueError:
                raise ValueError(
                    f"Invalid tenant_id: {tenant_id}. Must be a valid UUID."
                )

            # Ensure a knowledge base exists for this tenant
            kb_id = self._get_or_create_knowledge_base(tenant_id, knowledge_base_id)

            # PII redaction
            try:
                from ai_core.services.redactor import Redactor

                content = Redactor().sanitize(content)
            except Exception:
                pass
            # Optional encryption at rest
            raw_content = content
            try:
                from shared.crypto.crypto_service import crypto_service

                enc = crypto_service.encrypt(tenant_id, raw_content.encode("utf-8"))
                doc = Document(
                    title=title,
                    content="",
                    knowledge_base_id=kb_id,
                    status="PROCESSING",
                )
                setattr(doc, "raw_encrypted", enc)
                setattr(doc, "enc_ver", 1)
            except Exception:
                # Fallback to plaintext content if crypto not configured
                doc = Document(
                    title=title,
                    content=content,
                    knowledge_base_id=kb_id,
                    status="PROCESSING",
                )
            # Merge provided document-level metadata (RBAC, etc.)
            try:
                if isinstance(doc_meta, dict) and doc_meta:
                    base_meta = getattr(doc, "meta", {}) or {}
                    # Normalize keys
                    if "access" in doc_meta and isinstance(doc_meta["access"], str):
                        base_meta["access"] = doc_meta["access"]
                    if "owner_user_id" in doc_meta and doc_meta["owner_user_id"]:
                        base_meta["owner_user_id"] = str(doc_meta["owner_user_id"])
                    if "allowed_user_ids" in doc_meta and isinstance(
                        doc_meta["allowed_user_ids"], list
                    ):
                        base_meta["allowed_user_ids"] = [
                            str(x) for x in doc_meta["allowed_user_ids"]
                        ]
                    setattr(doc, "meta", base_meta)
            except Exception:
                pass
            self.db.add(doc)
            # Flush to assign PKs without committing the transaction
            self.db.flush()

            # AI-driven ingest plan and chunking
            plan = self.plan_ingest(title, content[:8000])
            chunk_pairs = self._build_chunks_with_metadata(raw_content)
            chunks = [t for (t, _m) in chunk_pairs]
            [m for (_t, m) in chunk_pairs]
            if not chunks:
                raise ValueError("No chunks could be created from the content")

            # Progress: chunking done
            if progress_job_id:
                try:
                    from shared.cache.redis import redis_cache

                    redis_cache.set_tenant_key(
                        tenant_id,
                        f"upload:job:{progress_job_id}",
                        {"phase": "embedding", "progress": 40},
                        ttl=3600,
                    )
                except Exception:
                    pass
            # Use planner directive for embedding path
            use_local_embeddings = bool(plan.get("use_local_embeddings"))
            # Decrypt if needed (already sanitized)
            try:
                if hasattr(doc, "raw_encrypted") and getattr(doc, "raw_encrypted"):
                    dec = crypto_service.decrypt(
                        tenant_id, getattr(doc, "raw_encrypted")
                    )
                    raw_content = dec.decode("utf-8", errors="ignore")
            except Exception:
                pass
            embeddings = self.embed(chunks, force_local=use_local_embeddings)

            # Store chunks
            qdrant_payload: List[Dict[str, Any]] = []
            for idx, ((chunk_text_val, meta_chunk), emb) in enumerate(
                zip(chunk_pairs, embeddings)
            ):
                # Ensure a concrete UUID is assigned before using the ID
                import uuid as _uuid

                chunk_id = _uuid.uuid4()
                # Merge auto-headline detection with chunk-derived meta
                chapter_meta = self._extract_chapter_info(chunk_text_val)
                merged_meta = dict(meta_chunk)
                for k, v in chapter_meta.items():
                    if v and k not in merged_meta:
                        merged_meta[k] = v
                kc = KnowledgeChunk(
                    id=chunk_id,
                    document_id=doc.id,
                    content=chunk_text_val,
                    chunk_index=idx,
                    embedding=emb,
                    meta=merged_meta or {},
                )
                self.db.add(kc)
                if progress_job_id and idx % 10 == 0:
                    try:
                        from shared.cache.redis import redis_cache

                        pct = 40 + int(50 * (idx + 1) / max(1, len(chunks)))
                        redis_cache.set_tenant_key(
                            tenant_id,
                            f"upload:job:{progress_job_id}",
                            {"phase": "storing", "progress": min(90, pct)},
                            ttl=3600,
                        )
                    except Exception:
                        pass
                # SQLAlchemy default UUID is assigned on instantiation; id is available before commit
                try:
                    # Extract ACL fields from document.meta
                    acl_access = None
                    owner_user_id = None
                    allowed_user_ids = None
                    try:
                        dmeta = getattr(doc, "meta", {}) or {}
                        acl_access = dmeta.get("access")
                        owner_user_id = dmeta.get("owner_user_id")
                        allowed_user_ids = dmeta.get("allowed_user_ids")
                    except Exception:
                        pass
                    payload = {
                        "id": str(chunk_id),
                        "embedding": emb,
                        "document_id": str(doc.id),
                        "document_title": doc.title,
                        "content": chunk_text_val,
                        "chunk_index": idx,
                        "chapter_num": merged_meta.get("chapter_num"),
                        "chapter_title": merged_meta.get("chapter_title"),
                        "page": merged_meta.get("page"),
                        "metadata": merged_meta or {},
                    }
                    # Attach ACL for vector-time filtering
                    if isinstance(acl_access, str):
                        payload["visibility"] = acl_access
                    if owner_user_id:
                        payload["owner_user_id"] = str(owner_user_id)
                    if isinstance(allowed_user_ids, list):
                        payload["allowed_user_ids"] = [str(x) for x in allowed_user_ids]
                    qdrant_payload.append(payload)
                except Exception:
                    pass
            doc.status = "INDEXED"
            doc.chunk_count = len(chunks)
            self.db.add(doc)
            self.db.flush()

            # Best-effort: upsert vectors to Qdrant when dimensions match (OpenAI = 1536)
            try:
                if qdrant_payload and isinstance(
                    qdrant_payload[0].get("embedding"), list
                ):
                    dim = (
                        len(qdrant_payload[0]["embedding"])
                        if qdrant_payload[0].get("embedding")
                        else 0
                    )
                    if dim in (1536, 3072, 1024):
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            # Collection may already exist or service may be unavailable
                            pass
                        try:
                            qdrant_service.upsert_knowledge_chunks(
                                tenant_id, qdrant_payload
                            )
                        except Exception as e:
                            logging.getLogger(__name__).warning(
                                f"Qdrant upsert skipped: {e}"
                            )
                            # enqueue for async retry (best-effort)
                            try:
                                from shared.queue.retry_queue import retry_queue

                                retry_queue.enqueue(
                                    job_type="qdrant_upsert",
                                    tenant_id=tenant_id,
                                    payload={"chunks": qdrant_payload},
                                    last_error=str(e),
                                )
                            except Exception:
                                pass
                    else:
                        logging.getLogger(__name__).info(
                            "Skipping Qdrant upsert due to embedding dimension mismatch"
                        )
            except Exception:
                # Never fail ingestion due to vector store issues
                pass

            # Write metadata.json for downstream processing
            metadata: Dict[str, Any] = {
                "tenant_id": tenant_id,
                "document_id": str(doc.id),
                "knowledge_base_id": kb_id,
                "title": title,
                "chunk_count": doc.chunk_count,
                "status": doc.status,
            }
            base_path = os.getenv(
                "DOCUMENT_STORAGE_PATH", os.path.join(os.getcwd(), "storage")
            )
            try:
                write_metadata(base_path, tenant_id, str(doc.id), metadata)
            except Exception as e:
                # Log but don't fail if metadata write fails
                import logging

                logger = logging.getLogger(__name__)
                logger.warning(f"Failed to write metadata: {e}")

            if progress_job_id:
                try:
                    from shared.cache.redis import redis_cache

                    redis_cache.set_tenant_key(
                        tenant_id,
                        f"upload:job:{progress_job_id}",
                        {"phase": "done", "progress": 100},
                        ttl=3600,
                    )
                except Exception:
                    pass
            # Commit the DB transaction only after all DB writes succeed
            trans.commit()
            try:
                from shared.metrics.ingestion_metrics import ingestion_metrics

                ingestion_metrics.inc_success(tenant_id, 1)
            except Exception:
                pass
            return str(doc.id), len(chunks)
        except Exception as e:
            # If there's an error, rollback the transaction
            try:
                trans.rollback()
            except Exception:
                pass
            try:
                from shared.metrics.ingestion_metrics import ingestion_metrics

                ingestion_metrics.inc_failure(locals().get("tenant_id", "global"), 1)
            except Exception:
                pass
            raise

    # ----------------------------
    # New modular pandas-based ingestion helpers
    # ----------------------------
    @staticmethod
    def load_file_to_dataframes(filename: str, data: bytes) -> Dict[str, pd.DataFrame]:
        """Detect file type, load into one or more pandas DataFrames.
        Returns a mapping of sheet_name -> DataFrame (for CSV, a single entry).
        Tries multiple header depths for hierarchical headers.
        """
        name = filename.lower()
        dfs: Dict[str, pd.DataFrame] = {}
        try:
            if name.endswith(".csv"):
                # Try common encodings and header depths
                text_variants = []
                for enc in ["utf-8-sig", "utf-8", "latin-1"]:
                    try:
                        text_variants.append(data.decode(enc))
                        break
                    except Exception:
                        continue
                raw = (
                    text_variants[0]
                    if text_variants
                    else data.decode("utf-8", errors="ignore")
                )
                for header_depth in [None, [0], [0, 1], [0, 1, 2]]:
                    try:
                        df = (
                            pd.read_csv(pd.io.common.StringIO(raw), header=header_depth)
                            if header_depth is not None
                            else pd.read_csv(pd.io.common.StringIO(raw))
                        )
                        if df is not None and df.shape[0] > 0:
                            dfs["Sheet1"] = df
                            break
                    except Exception:
                        continue
                if not dfs:
                    # final fallback
                    dfs["Sheet1"] = pd.read_csv(pd.io.common.StringIO(raw), header=0)
            elif name.endswith(".xlsx"):
                buf = io.BytesIO(data)
                # Try multiple header depths per sheet
                xls = pd.ExcelFile(buf, engine="openpyxl")
                for sheet in xls.sheet_names:
                    df: pd.DataFrame | None = None
                    for header_depth in [[0, 1, 2], [0, 1], [0]]:
                        try:
                            df_try = pd.read_excel(
                                xls, sheet_name=sheet, header=header_depth
                            )
                            if df_try is not None and df_try.shape[0] > 0:
                                df = df_try
                                break
                        except Exception:
                            continue
                    if df is None:
                        df = pd.read_excel(xls, sheet_name=sheet)
                    dfs[sheet] = df
            else:
                raise ValueError(
                    "Unsupported tabular file type; expected .csv or .xlsx"
                )
            # Log basic info
            logging.getLogger(__name__).info(
                f"Loaded file '{filename}' into {len(dfs)} DataFrame(s): {list(dfs.keys())}"
            )
            return dfs
        except Exception as e:
            logging.getLogger(__name__).error(
                f"Failed to load '{filename}' into pandas: {e}"
            )
            raise

    @staticmethod
    def normalize_headers(df: pd.DataFrame) -> pd.DataFrame:
        """Flatten multi-level headers and strip empty parts."""
        try:
            if isinstance(df.columns, pd.MultiIndex):
                df = df.copy()
                df.columns = [
                    " | ".join([str(c) for c in col if str(c) != "nan"]).strip()
                    for col in df.columns.values
                ]
            else:
                df = df.copy()
                df.columns = [str(c).strip() for c in df.columns]
        except Exception:
            # Best-effort normalization
            df = df.copy()
            df.columns = [str(c) for c in df.columns]
        return df

    @staticmethod
    def _is_mostly_numeric_or_nan(series: pd.Series) -> bool:
        try:
            s = pd.to_numeric(series, errors="coerce")
            frac_num_or_nan = float(s.notna().sum()) / max(1, len(s))
            # Consider numeric if original non-nulls mostly converted to numeric
            return frac_num_or_nan >= 0.9
        except Exception:
            return False

    @staticmethod
    def _should_summarize(df: pd.DataFrame) -> bool:
        try:
            non_obj_cols = [
                c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])
            ]
            frac_numeric = float(len(non_obj_cols)) / max(1, len(df.columns))
            return (frac_numeric >= 0.7) and (len(df.columns) < 10)
        except Exception:
            return False

    @staticmethod
    def _summarize_numeric_df(df: pd.DataFrame, title: str) -> str:
        try:
            desc = df.describe(include=[np.number]).to_dict()
            lines = [f"Summary for {title}:"]
            for col, stats in desc.items():
                if not isinstance(stats, dict):
                    continue
                mn = stats.get("min")
                mx = stats.get("max")
                mean = stats.get("mean")
                lines.append(
                    f"- {col}: min={mn}, max={mx}, mean={round(mean, 3) if mean is not None else mean}"
                )
            # Simple outlier detection: z-score > 3 (approx.)
            try:
                z = (
                    df[[c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]]
                    - df.mean(numeric_only=True)
                ) / df.std(numeric_only=True)
                outlier_counts = (np.abs(z) > 3).sum().to_dict()
                for col, cnt in outlier_counts.items():
                    if int(cnt) > 0:
                        lines.append(f"- Outliers detected in {col}: {int(cnt)} rows")
            except Exception:
                pass
            return "\n".join(lines)
        except Exception:
            return f"Summary for {title}: numeric overview unavailable."

    @classmethod
    def df_to_semantic_docs(
        cls, df: pd.DataFrame, filename: str, sheet_name: str | None = None
    ) -> List[Tuple[str, Dict[str, Any]]]:
        """Convert a DataFrame into row-wise semantic mini-documents.
        If the DataFrame is predominantly numeric with few columns, return a single summary doc.
        Returns list of (text, metadata) where metadata includes source_file and row_index.
        """
        docs: List[Tuple[str, Dict[str, Any]]] = []
        if df is None or df.shape[0] == 0:
            return docs
        df = cls.normalize_headers(df)
        # Decide to summarize vs row-wise
        if cls._should_summarize(df):
            text = cls._summarize_numeric_df(df, title=sheet_name or filename)
            meta = {"source_file": filename, "sheet": sheet_name, "summary": True}
            return [(text, meta)]
        # Identify fields to ignore (mostly numeric/NaN)
        ignore_cols = set()
        for c in df.columns:
            try:
                if cls._is_mostly_numeric_or_nan(df[c]):
                    ignore_cols.add(c)
            except Exception:
                continue
        # Build per-row docs and include key-value map in metadata
        for idx, row in df.iterrows():
            parts: List[str] = [f"Record {idx}:"]
            row_map: Dict[str, Any] = {}
            for c in df.columns:
                if c in ignore_cols:
                    continue
                val = row.get(c)
                if pd.isna(val) or (isinstance(val, str) and not val.strip()):
                    continue
                # Include in semantic text
                parts.append(f"{c}: {val}")
                # Always include in structured metadata map (even if numeric) for precise retrieval later
                try:
                    row_map[str(c)] = "" if pd.isna(val) else str(val)
                except Exception:
                    row_map[str(c)] = str(val)
            text = "\n".join(parts)
            if text.strip() and len(parts) > 1:
                meta = {"source_file": filename, "row_index": int(idx), "row": row_map}
                if sheet_name:
                    meta["sheet"] = sheet_name
                docs.append((text, meta))
        return docs

    def process_pandas_and_store(
        self,
        tenant_id: str,
        title: str,
        filename: str,
        data: bytes,
        knowledge_base_id: str,
        progress_job_id: str | None = None,
    ) -> Tuple[str, int]:
        """End-to-end: load file to DataFrame(s), convert to semantic docs, embed and store one chunk per doc.
        Returns (document_id_of_first, total_chunks).
        """
        # Plan enforcement: file size and max_docs
        try:
            from shared.database.models import KnowledgeBase, Document as _Doc

            from sqlalchemy import func as _f
            s_bytes = len(data or b"")
            # Determine plan label
            t = self.db.query(Tenant).filter(Tenant.id == tenant_id).first()
            plan_label = resolve_plan_label(getattr(t, "subscription_tier", None))
            plan = get_plan(plan_label)
            max_size = int(plan.get("limits", {}).get("max_file_size") or 0)
            if max_size and s_bytes > max_size:
                raise ValueError("File exceeds plan file size limit. Upgrade plan.")
            # Count docs for tenant
            cnt = (
                self.db.query(_Doc)
                .join(KnowledgeBase, _Doc.knowledge_base_id == KnowledgeBase.id)
                .filter(KnowledgeBase.tenant_id == tenant_id)
                .count()
            )
            max_docs = int(plan.get("limits", {}).get("max_docs") or 0)
            if max_docs and cnt >= max_docs:
                raise ValueError("Maximum documents reached for your plan. Upgrade to add more.")
        except Exception as _e:
            # proceed if enforcement queries fail; errors will surface elsewhere
            pass
        # Load multiple DataFrames (sheets)
        dfs = self.load_file_to_dataframes(filename, data)
        logger = logging.getLogger(__name__)
        logger.info(f"Ingesting '{filename}': sheets={list(dfs.keys())}")
        # Create/ensure KB
        kb_id = self._get_or_create_knowledge_base(tenant_id, knowledge_base_id)
        total_chunks = 0
        first_doc_id: str | None = None
        for sheet, df in dfs.items():
            # Convert to semantic documents
            docs = self.df_to_semantic_docs(df, filename=filename, sheet_name=sheet)
            if not docs:
                logger.info(f"No documents generated for sheet '{sheet}'")
                continue
            # Create parent Document row with a preview
            preview = "\n\n".join([d[0] for d in docs[:2]])
            parent = Document(
                title=f"{title} - {sheet}",
                content=preview,
                knowledge_base_id=kb_id,
                status="PROCESSING",
            )
            # Optionally store flattened columns on meta
            try:
                parent.meta = {"columns": list(df.columns)}
            except Exception:
                pass
            self.db.add(parent)
            self.db.commit()
            self.db.refresh(parent)
            # Persist schema fields in cache for query expansion
            try:
                from shared.cache.redis import redis_cache

                cols_norm = [str(c).strip() for c in df.columns]
                existing = redis_cache.get_tenant_key(tenant_id, "schema:fields")
                merged = []
                if isinstance(existing, list):
                    merged = list(dict.fromkeys([*existing, *cols_norm]))
                else:
                    merged = cols_norm
                redis_cache.set_tenant_key(
                    tenant_id, "schema:fields", merged, ttl=24 * 3600
                )
            except Exception:
                pass
            # Progress update
            if progress_job_id:
                try:
                    from shared.cache.redis import redis_cache

                    redis_cache.set_tenant_key(
                        tenant_id,
                        f"upload:job:{progress_job_id}",
                        {"phase": "embedding", "progress": 40},
                        ttl=3600,
                    )
                except Exception:
                    pass
            # Embed docs (use local for robustness by default for tabular)
            # PII redaction for tabular docs
            try:
                from ai_core.services.redactor import Redactor

                red = Redactor()
                texts = [red.sanitize(t) for (t, _m) in docs]
            except Exception:
                texts = [t for (t, _m) in docs]
            metas = [m for (_t, m) in docs]
            # Prefer remote embeddings when available (3072/1536 dims) for Qdrant upsert; fallback to local 256-d
            use_local = False if os.getenv("OPENAI_API_KEY") else True
            embeddings = self.embed(texts, force_local=use_local)
            # Store rows as KnowledgeChunks
            qdrant_payload: List[Dict[str, Any]] = []
            for idx, (t, m, emb) in enumerate(zip(texts, metas, embeddings)):
                import uuid as _uuid

                chunk_id = _uuid.uuid4()
                # Encrypt chunk content; store preview only in plaintext
                try:
                    enc = crypto_service.encrypt(tenant_id, t.encode("utf-8"))
                    kc = KnowledgeChunk(
                        id=chunk_id,
                        document_id=parent.id,
                        content=t[:160],
                        content_encrypted=enc,
                        enc_ver=1,
                        chunk_index=idx,
                        embedding=emb,
                        meta=m,
                    )
                except Exception:
                    kc = KnowledgeChunk(
                        id=chunk_id,
                        document_id=parent.id,
                        content=t,
                        chunk_index=idx,
                        embedding=emb,
                        meta=m,
                    )
                self.db.add(kc)
                total_chunks += 1
                # Prepare optional vector payload
                try:
                    qdrant_payload.append(
                        {
                            "id": str(chunk_id),
                            "embedding": emb,
                            "document_id": str(parent.id),
                            "document_title": parent.title,
                            "content": t[:160],
                            "chunk_index": idx,
                            "chapter_num": None,
                            "chapter_title": None,
                            "page": None,
                        }
                    )
                except Exception:
                    pass
            parent.status = "INDEXED"
            parent.chunk_count = len(texts)
            self.db.add(parent)
            self.db.commit()
            if first_doc_id is None:
                first_doc_id = str(parent.id)
            # Log sample preview
            logger.info(
                f"Parsed rows for '{sheet}': {len(texts)}; sample=\n{texts[0][:400]}"
            )
            # Best-effort upsert to Qdrant if dims match
            try:
                if qdrant_payload and isinstance(
                    qdrant_payload[0].get("embedding"), list
                ):
                    dim = (
                        len(qdrant_payload[0]["embedding"])
                        if qdrant_payload[0].get("embedding")
                        else 0
                    )
                    if dim in (1536, 3072, 1024):
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            pass
                        try:
                            qdrant_service.upsert_knowledge_chunks(
                                tenant_id, qdrant_payload
                            )
                        except Exception as e:
                            logger.warning(f"Qdrant upsert skipped: {e}")
            except Exception:
                pass

            # Field–Value Semantic Chunking: upsert per-cell vectors (1536-dim)
            try:
                fv_texts: List[str] = []
                fv_meta: List[Dict[str, Any]] = []

                def _norm(s: str) -> str:
                    return re.sub(r"[^a-z0-9]+", "_", str(s).strip().lower()).strip("_")

                for ridx, row in df.iterrows():
                    record_sig = f"{title or ''}|{sheet or ''}|{ridx}"
                    for col in list(df.columns):
                        try:
                            raw_val = row.get(col)
                        except Exception:
                            try:
                                raw_val = row[col]
                            except Exception:
                                raw_val = None
                        if pd.isna(raw_val) or str(raw_val).strip() == "":
                            continue
                        field_display = str(col).strip()
                        field_name = _norm(field_display)
                        # Robust value normalization for non-text fields
                        try:
                            if isinstance(raw_val, (int, float)):
                                value_raw = f"{raw_val}"
                            else:
                                value_raw = str(raw_val)
                        except Exception:
                            value_raw = str(raw_val) if raw_val is not None else ""
                        value_raw = value_raw.strip()
                        # Add simple hierarchical context (sheet and filename)
                        context_bits = []
                        if sheet:
                            context_bits.append(f"Sheet: {sheet}")
                        if title:
                            context_bits.append(f"File: {title}")
                        ctx_str = " | ".join(context_bits)
                        text = (
                            f"Field: {field_display} | Value: {value_raw} | Record: {int(ridx)+1}"
                            + (f" | {ctx_str}" if ctx_str else "")
                        )
                        try:
                            from ai_core.services.redactor import Redactor

                            text = Redactor().sanitize(text)
                        except Exception:
                            pass
                        fv_texts.append(text)
                        # Basic importance heuristic: de-emphasize id-like fields
                        importance = 0.8
                        if re.search(r"\b(id|uuid|ssn|number|no)\b", field_name):
                            importance = 0.3
                        elif re.search(
                            r"\b(date|dob|salary|amount|manager|department|title|position)\b",
                            field_name,
                        ):
                            importance = 1.2
                        fv_meta.append(
                            {
                                "row_index": int(ridx),
                                "field_name": field_name,
                                "field_display": field_display,
                                "value_raw": value_raw,
                                "value_norm": _norm(value_raw),
                                "sheet": sheet,
                                "source_file": title,
                                "record_sig": record_sig,
                                "document_id": str(parent.id),
                                "content": text,
                                "importance": float(importance),
                            }
                        )
                if fv_texts:
                    emb = self.embed(fv_texts, force_local=False)
                    ready: List[Dict[str, Any]] = []
                    if emb and all(
                        isinstance(v, list) and len(v) in (1536, 3072, 1024)
                        for v in emb
                    ):
                        import uuid as _uuid

                        for m, vec in zip(fv_meta, emb):
                            fid = str(
                                _uuid.uuid5(
                                    _uuid.NAMESPACE_URL,
                                    f"{tenant_id}:fv:{m['document_id']}:{m['row_index']}:{m['field_name']}",
                                )
                            )
                            item = {
                                "id": fid,
                                "embedding": vec,
                            }
                            item.update(m)
                            ready.append(item)
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            pass
                        try:
                            qdrant_service.upsert_field_values(tenant_id, ready)
                        except Exception as e:
                            logger.warning(f"Qdrant field_values upsert skipped: {e}")
            except Exception as e:
                logger.warning(f"Field–Value semantic chunking skipped: {e}")

            # Upsert schema field embeddings (one-time per sheet) into schema collection
            try:
                # Build simple descriptions per column
                col_texts: List[str] = []
                col_names: List[str] = []
                col_aliases: List[List[str]] = []
                col_tokens: List[List[str]] = []
                for c in list(df.columns):
                    cname = str(c).strip()
                    if not cname:
                        continue
                    desc = f"Field: {cname}. Meaning: {cname.replace('_', ' ')}."
                    # naive alias/token generation for bootstrap
                    tokens = [t for t in re.split(r"[^a-z0-9]+", cname.lower()) if t]
                    aliases = list(
                        dict.fromkeys(
                            [
                                cname,
                                cname.replace("_", " "),
                                " ".join(tokens),
                            ]
                        )
                    )
                    col_texts.append(desc)
                    col_names.append(cname)
                    col_aliases.append(aliases)
                    col_tokens.append(tokens)
                if col_texts:
                    # Use remote embeddings to match Qdrant 1536 dim
                    schema_embs = self.embed(col_texts, force_local=False)
                    # Validate dimensions
                    ready: List[Dict[str, Any]] = []
                    for i, (cname, emb) in enumerate(zip(col_names, schema_embs)):
                        if isinstance(emb, list) and len(emb) in (1536, 3072, 1024):
                            import uuid as _uuid

                            # Deterministic ID per tenant+field
                            fid = str(
                                _uuid.uuid5(
                                    _uuid.NAMESPACE_URL, f"{tenant_id}:schema:{cname}"
                                )
                            )
                            ready.append(
                                {
                                    "id": fid,
                                    "embedding": emb,
                                    "field_name": cname,
                                    "description": f"{cname}",
                                    "aliases": col_aliases[i],
                                    "tokens": col_tokens[i],
                                }
                            )
                    if ready:
                        # Cleanup: delete stale schema fields no longer present
                        try:
                            existing = qdrant_service.list_schema_fields(
                                tenant_id, limit=2000
                            )
                            existing_names = set()
                            id_by_name = {}
                            for p in existing:
                                pl = p.get("payload") or {}
                                nm = pl.get("field_name")
                                if isinstance(nm, str):
                                    existing_names.add(nm)
                                    id_by_name[nm] = p.get("id")
                            current_names = set(col_names)
                            stale_names = list(existing_names - current_names)
                            stale_ids = [
                                id_by_name[n] for n in stale_names if n in id_by_name
                            ]
                            if stale_ids:
                                qdrant_service.delete_schema_fields_by_ids(
                                    tenant_id, stale_ids
                                )
                        except Exception:
                            pass
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            pass
                        qdrant_service.upsert_schema_fields(tenant_id, ready)
            except Exception as e:
                logger.warning(f"Schema embedding upsert skipped: {e}")
        if not first_doc_id:
            raise ValueError("No data parsed from file")
        # Final progress
        if progress_job_id:
            try:
                from shared.cache.redis import redis_cache

                redis_cache.set_tenant_key(
                    tenant_id,
                    f"upload:job:{progress_job_id}",
                    {"phase": "done", "progress": 100},
                    ttl=3600,
                )
            except Exception:
                pass
        return first_doc_id, total_chunks

    def _get_or_create_knowledge_base(self, tenant_id: str, provided_kb_id: str) -> str:
        import uuid

        # Validate tenant_id
        try:
            tenant_uuid = uuid.UUID(tenant_id)
        except ValueError:
            raise ValueError(f"Invalid tenant_id format: {tenant_id}")

        # Ensure tenant exists; if missing, create a default record (dev-friendly)
        try:
            existing_tenant = self.db.get(Tenant, tenant_uuid)
            if not existing_tenant:
                # Create a minimal tenant to satisfy FK; name/domain deterministic for the given UUID
                t = Tenant(
                    id=tenant_uuid, name="Seeded Tenant", domain="seeded", settings={}
                )
                self.db.add(t)
                self.db.commit()
        except Exception:
            # Best-effort; if this fails, the subsequent KB creation will surface the error
            self.db.rollback()

        # If a valid kb id is provided and exists, use it
        if provided_kb_id and provided_kb_id != "00000000-0000-0000-0000-000000000000":
            try:
                kb_uuid = uuid.UUID(provided_kb_id)
                kb = self.db.get(KnowledgeBase, kb_uuid)
                if kb:
                    return str(kb.id)
            except ValueError:
                # Invalid UUID format, skip
                pass

        # Try to find default KB for tenant
        kb = (
            self.db.query(KnowledgeBase)
            .filter(KnowledgeBase.tenant_id == tenant_uuid)
            .order_by(KnowledgeBase.created_at.asc())
            .first()
        )
        if kb:
            return str(kb.id)

        # Create a new knowledge base for this tenant
        new_kb = KnowledgeBase(
            tenant_id=tenant_uuid, name="Default", status="ACTIVE", document_count=0
        )
        self.db.add(new_kb)
        self.db.commit()
        self.db.refresh(new_kb)
        return str(new_kb.id)

    def extract_text_from_file(self, filename: str, data: bytes) -> str:
        name = filename.lower()
        if name.endswith(".txt") or name.endswith(".csv"):
            try:
                return data.decode("utf-8")
            except Exception:
                return data.decode("latin-1", errors="ignore")
        if name.endswith(".docx"):
            buf = io.BytesIO(data)
            d = DocxDocument(buf)
            return "\n".join(p.text for p in d.paragraphs)
        if name.endswith(".pptx"):
            buf = io.BytesIO(data)
            prs = Presentation(buf)
            texts: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        if name.endswith(".xlsx"):
            buf = io.BytesIO(data)
            wb = load_workbook(buf, data_only=True)
            texts: List[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    texts.append("\t".join("" if v is None else str(v) for v in row))
            return "\n".join(texts)
        # PDF handled by PyPDF2 with OCR fallback
        if name.endswith(".pdf"):
            try:
                from PyPDF2 import PdfReader

                buf = io.BytesIO(data)
                reader = PdfReader(buf)
                texts: List[str] = []
                for i, page in enumerate(reader.pages, start=1):
                    extracted = page.extract_text() or ""
                    if not extracted.strip():
                        # OCR fallback via pdfplumber + pytesseract
                        try:
                            import pdfplumber
                            import pytesseract
                            from PIL import Image

                            with pdfplumber.open(io.BytesIO(data)) as pdf:
                                if 0 <= (i - 1) < len(pdf.pages):
                                    p = pdf.pages[i - 1]
                                    im = p.to_image(resolution=200).original
                                    # Ensure PIL Image
                                    if not isinstance(im, Image.Image):
                                        im = Image.fromarray(im)
                                    ocr_text = pytesseract.image_to_string(im)
                                    extracted = ocr_text or ""
                        except Exception:
                            pass
                    texts.append(f"[[PAGE:{i}]]\n" + extracted)
                return "\n".join(texts)
            except Exception:
                pass
        # Fallback raw decode and OCR for common image formats
        try:
            return data.decode("utf-8")
        except Exception:
            # Try image OCR if this appears to be an image
            try:
                import imghdr

                kind = imghdr.what(None, h=data)
                if kind in ("png", "jpeg", "jpg", "bmp", "tiff"):
                    from PIL import Image
                    import pytesseract

                    img = Image.open(io.BytesIO(data))
                    return pytesseract.image_to_string(img)
            except Exception:
                pass
            return data.decode("latin-1", errors="ignore")

    def extract_rows_from_file(self, filename: str, data: bytes) -> List[str]:
        name = filename.lower()
        rows: List[str] = []
        if name.endswith(".csv"):
            try:
                text = data.decode("utf-8-sig", errors="ignore")
            except Exception:
                text = data.decode("latin-1", errors="ignore")
            reader = csv.reader(io.StringIO(text))
            # Re-serialize each row via csv.writer to preserve quoting and commas inside fields
            for r in reader:
                buf = io.StringIO()
                writer = csv.writer(buf)
                writer.writerow(r)
                rows.append(buf.getvalue().strip("\n"))
        elif name.endswith(".xlsx"):
            buf = io.BytesIO(data)
            wb = load_workbook(buf, data_only=True)
            for ws in wb.worksheets:
                for r in ws.iter_rows(values_only=True):
                    # Use csv.writer to serialize each row consistently
                    buf_row = io.StringIO()
                    writer = csv.writer(buf_row)
                    writer.writerow(["" if v is None else v for v in r])
                    rows.append(buf_row.getvalue().strip("\n"))
        return rows

    def extract_rows_by_sheet(self, filename: str, data: bytes) -> Dict[str, List[str]]:
        name = filename.lower()
        result: Dict[str, List[str]] = {}
        if name.endswith(".xlsx"):
            buf = io.BytesIO(data)
            wb = load_workbook(buf, data_only=True)
            for ws in wb.worksheets:
                rows: List[str] = []
                for r in ws.iter_rows(values_only=True):
                    buf_row = io.StringIO()
                    writer = csv.writer(buf_row)
                    writer.writerow(["" if v is None else v for v in r])
                    rows.append(buf_row.getvalue().strip("\n"))
                result[ws.title] = rows
        return result

    def process_rows_and_store(
        self,
        tenant_id: str,
        title: str,
        rows: List[str],
        knowledge_base_id: str,
        progress_job_id: str | None = None,
        sheet_name: str | None = None,
    ) -> Tuple[str, int]:
        try:
            # Validate tenant_id
            import uuid

            try:
                uuid.UUID(tenant_id)
            except ValueError:
                raise ValueError(
                    f"Invalid tenant_id: {tenant_id}. Must be a valid UUID."
                )

            if not rows:
                raise ValueError("No rows provided to process")

            kb_id = self._get_or_create_knowledge_base(tenant_id, knowledge_base_id)
            preview = "\n".join(rows[:5]) + ("\n..." if len(rows) > 5 else "")
            doc = Document(
                title=title,
                content=preview,
                knowledge_base_id=kb_id,
                status="PROCESSING",
            )

            # Capture header columns if present (first row)
            if rows:
                try:
                    header_reader = csv.reader(io.StringIO(rows[0]))
                    header = next(header_reader)
                    doc.meta = {"columns": [h.strip().lower() for h in header]}
                    # If first row is header, skip it for chunk storage
                    data_rows = rows[1:]
                except Exception:
                    data_rows = rows
            else:
                data_rows = rows

            if not data_rows:
                raise ValueError("No data rows found after header extraction")

            self.db.add(doc)
            self.db.commit()
            self.db.refresh(doc)

            if progress_job_id:
                try:
                    from shared.cache.redis import redis_cache

                    redis_cache.set_tenant_key(
                        tenant_id,
                        f"upload:job:{progress_job_id}",
                        {"phase": "embedding", "progress": 40},
                        ttl=3600,
                    )
                except Exception:
                    pass
            # Redact PII in raw row strings
            try:
                from ai_core.services.redactor import Redactor

                red = Redactor()
                data_rows = [red.sanitize(r) for r in data_rows]
            except Exception:
                pass
            # For tabular rows, use local deterministic embeddings to avoid API payload rejections
            use_local = False if os.getenv("OPENAI_API_KEY") else True
            embeddings = self.embed(data_rows, force_local=use_local)
            qdrant_payload: List[Dict[str, Any]] = []
            for idx, (row_text, emb) in enumerate(zip(data_rows, embeddings)):
                import uuid as _uuid

                chunk_id = _uuid.uuid4()
                # Attach parsed row and sheet metadata for deterministic aggregation
                meta_row: Dict[str, Any] = {}
                try:
                    header = (
                        (doc.meta or {}).get("columns")
                        if isinstance(doc.meta, dict)
                        else None
                    )
                    if isinstance(header, list):
                        reader = csv.reader(io.StringIO(row_text))
                        row_vals = next(reader)
                        row_map = {}
                        for i, col in enumerate(header):
                            if i < len(row_vals):
                                row_map[str(col).strip().lower()] = row_vals[i]
                        meta_row["row"] = row_map
                except Exception:
                    pass
                if sheet_name:
                    meta_row["sheet"] = sheet_name
                # Encrypt and store preview
                try:
                    enc = crypto_service.encrypt(tenant_id, row_text.encode("utf-8"))
                    kc = KnowledgeChunk(
                        id=chunk_id,
                        document_id=doc.id,
                        content=row_text[:160],
                        content_encrypted=enc,
                        enc_ver=1,
                        chunk_index=idx,
                        embedding=emb,
                        meta=meta_row,
                    )
                except Exception:
                    kc = KnowledgeChunk(
                        id=chunk_id,
                        document_id=doc.id,
                        content=row_text,
                        chunk_index=idx,
                        embedding=emb,
                        meta=meta_row,
                    )
                self.db.add(kc)
                if progress_job_id and idx % 50 == 0:
                    try:
                        from shared.cache.redis import redis_cache

                        pct = 40 + int(50 * (idx + 1) / max(1, len(data_rows)))
                        redis_cache.set_tenant_key(
                            tenant_id,
                            f"upload:job:{progress_job_id}",
                            {"phase": "storing", "progress": min(90, pct)},
                            ttl=3600,
                        )
                    except Exception:
                        pass
                try:
                    qdrant_payload.append(
                        {
                            "id": str(chunk_id),
                            "embedding": emb,
                            "document_id": str(doc.id),
                            "document_title": doc.title,
                            "content": row_text,
                            "chunk_index": idx,
                            "chapter_num": None,
                            "chapter_title": None,
                            "metadata": meta_row or {},
                        }
                    )
                except Exception:
                    pass
            doc.status = "INDEXED"
            doc.chunk_count = len(data_rows)
            self.db.add(doc)
            self.db.commit()

            # Best-effort: upsert to Qdrant when dimensions match expected size
            try:
                if qdrant_payload and isinstance(
                    qdrant_payload[0].get("embedding"), list
                ):
                    dim = (
                        len(qdrant_payload[0]["embedding"])
                        if qdrant_payload[0].get("embedding")
                        else 0
                    )
                    if dim in (1536, 3072, 1024):
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            pass
                        try:
                            qdrant_service.upsert_knowledge_chunks(
                                tenant_id, qdrant_payload
                            )
                        except Exception as e:
                            logging.getLogger(__name__).warning(
                                f"Qdrant upsert skipped: {e}"
                            )
                    else:
                        logging.getLogger(__name__).info(
                            "Skipping Qdrant upsert due to embedding dimension mismatch"
                        )
            except Exception:
                pass

            # Upsert schema field embeddings for CSV header if present
            try:
                header = (
                    (doc.meta or {}).get("columns")
                    if isinstance(doc.meta, dict)
                    else None
                )
                if isinstance(header, list) and header:
                    col_texts: List[str] = []
                    col_names: List[str] = []
                    for c in header:
                        cname = str(c).strip()
                        if not cname:
                            continue
                        desc = f"Field: {cname}. Meaning: {cname.replace('_', ' ')}."
                        col_texts.append(desc)
                        col_names.append(cname)
                    if col_texts:
                        schema_embs = self.embed(col_texts, force_local=False)
                        ready: List[Dict[str, Any]] = []
                        for cname, emb in zip(col_names, schema_embs):
                            if isinstance(emb, list) and len(emb) in (1536, 3072, 1024):
                                import uuid as _uuid

                                fid = str(
                                    _uuid.uuid5(
                                        _uuid.NAMESPACE_URL,
                                        f"{tenant_id}:schema:{cname}",
                                    )
                                )
                                ready.append(
                                    {
                                        "id": fid,
                                        "embedding": emb,
                                        "field_name": cname,
                                        "description": f"{cname}",
                                    }
                                )
                        if ready:
                            try:
                                qdrant_service.create_collection()
                            except Exception:
                                pass
                            qdrant_service.upsert_schema_fields(tenant_id, ready)
            except Exception:
                pass

            metadata: Dict[str, Any] = {
                "tenant_id": tenant_id,
                "document_id": str(doc.id),
                "knowledge_base_id": kb_id,
                "title": title,
                "chunk_count": doc.chunk_count,
                "status": doc.status,
            }
            base_path = os.getenv(
                "DOCUMENT_STORAGE_PATH", os.path.join(os.getcwd(), "storage")
            )
            try:
                write_metadata(base_path, tenant_id, str(doc.id), metadata)
            except Exception as e:
                # Log but don't fail if metadata write fails
                import logging

                logger = logging.getLogger(__name__)
                logger.warning(f"Failed to write metadata: {e}")

            if progress_job_id:
                try:
                    from shared.cache.redis import redis_cache

                    redis_cache.set_tenant_key(
                        tenant_id,
                        f"upload:job:{progress_job_id}",
                        {"phase": "done", "progress": 100},
                        ttl=3600,
                    )
                except Exception:
                    pass
            return str(doc.id), len(rows)
        except Exception as e:
            # Rollback on error
            self.db.rollback()
            raise
